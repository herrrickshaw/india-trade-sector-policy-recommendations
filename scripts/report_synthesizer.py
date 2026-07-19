#!/usr/bin/env python3
"""
report_synthesizer.py — turn any bulletin's data JSON in this repo into a
standalone PDF, Word (.docx), and PowerPoint (.pptx) report.

Why this exists: every bulletin in this repo is a hand-designed HTML chart
plus a sourced data JSON. When the underlying figures get refreshed (a new
DPIIT quarter, a new PIB release, a rate change), rebuilding the *chart* is
a design task -- but producing a *shareable report* of the current numbers
shouldn't be. This script does that mechanically, from the JSON alone.

Two ways to get a PDF:
  --html <chart.html>   Renders the paired chart bulletin exactly as designed,
                         via headless Chrome's print-to-PDF. Pixel-faithful to
                         the published page. Requires Chrome/Chromium installed.
  (no --html)           Builds a plain HTML page from the JSON's own generic
                         structure (see below) and prints THAT to PDF instead.
                         Falls back to a lightweight styled page — not the
                         designed bulletin, but always available.

The Word and PowerPoint outputs are ALWAYS built from the JSON directly, by
walking its structure generically:
  - a dict becomes a heading, then each of its keys recurses one level deeper
  - a list of same-shaped dicts becomes a table (union of keys as columns)
  - a list of scalars becomes a bullet list
  - a scalar becomes a paragraph
This means the script does NOT need a template per bulletin. It keeps working
as new bulletins with new JSON shapes get added -- the only conventions it
leans on are the ones this repo already follows: a top-level "purpose" and
"retrieved" field if present are used as the report's subtitle/date.

USAGE
-----
  # Single bulletin, all three formats, matching the exact published chart:
  python3 report_synthesizer.py ../data/fdi_vs_pli_launch_2026-07-19.json \\
      --html ../charts/fdi_vs_pli_launch_correlation.html \\
      --formats pdf,docx,pptx --out-dir ../reports/

  # JSON only (no paired chart) -- docx/pptx always work, pdf uses the
  # generic-HTML fallback:
  python3 report_synthesizer.py ../data/sector_master_scorecard_2026-07-19.json \\
      --formats docx,pptx

  # Batch every bulletin listed in the repo manifest:
  python3 report_synthesizer.py --manifest ../data/repo_manifest_2026-07-19.json \\
      --formats pdf,docx,pptx --out-dir ../reports/

REQUIREMENTS
------------
  pip install -r requirements.txt   (python-docx, python-pptx; pypdf optional,
                                      used only to report page counts)
  A local Chrome/Chromium install for any --formats that include pdf.
"""

import argparse
import json
import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

# ---------------------------------------------------------------------------
# Small, repo-specific vocabulary so headings read naturally instead of as
# raw_json_keys. Extend this dict rather than special-casing key names below.
# ---------------------------------------------------------------------------
ACRONYMS = {
    "fdi": "FDI", "pli": "PLI", "hs": "HS", "ibef": "IBEF", "icra": "ICRA",
    "gnep": "GNEP", "rbi": "RBI", "mclr": "MCLR", "eblr": "EBLR", "walr": "WALR",
    "usd": "USD", "cbam": "CBAM", "dpiit": "DPIIT", "pib": "PIB", "niti": "NITI",
    "aayog": "Aayog", "ceta": "CETA", "sez": "SEZ", "iem": "IEM", "sia": "SIA",
    "gdp": "GDP", "cpi": "CPI", "wpi": "WPI",
}


def humanize(key):
    """snake_case_key -> Title Case Key, honoring ACRONYMS."""
    words = str(key).replace("-", "_").split("_")
    out = []
    for w in words:
        lw = w.lower()
        out.append(ACRONYMS.get(lw, w.capitalize() if w else w))
    return " ".join(out) if out else str(key)


def title_from_stem(stem):
    """Filename stem -> a readable default title, used only when no --manifest
    title override is supplied. Strips a trailing _YYYY_MM_DD date stamp
    (this repo's own naming convention) before humanizing."""
    import re
    stripped = re.sub(r"_\d{4}_\d{2}_\d{2}$", "", str(stem).replace("-", "_"))
    return humanize(stripped)


def stringify_cell(value, max_len=280):
    """Collapse an arbitrary JSON value into one table-cell-sized string."""
    if value is None:
        return ""
    if isinstance(value, (str, int, float, bool)):
        s = str(value)
    elif isinstance(value, list):
        if all(isinstance(v, (str, int, float, bool)) or v is None for v in value):
            s = "; ".join(stringify_cell(v, 80) for v in value)
        else:
            s = "; ".join(stringify_cell(v, 100) for v in value)
    elif isinstance(value, dict):
        s = "; ".join(f"{humanize(k)}: {stringify_cell(v, 60)}" for k, v in value.items())
    else:
        s = str(value)
    if len(s) > max_len:
        s = s[: max_len - 1].rstrip() + "…"
    return s


def is_list_of_dicts(value):
    return isinstance(value, list) and len(value) > 0 and all(isinstance(v, dict) for v in value)


def is_list_of_scalars(value):
    return isinstance(value, list) and len(value) > 0 and all(
        isinstance(v, (str, int, float, bool)) or v is None for v in value
    )


def table_columns(rows):
    """Union of keys across a list of dicts, in first-seen order."""
    cols = []
    seen = set()
    for row in rows:
        for k in row.keys():
            if k not in seen:
                seen.add(k)
                cols.append(k)
    return cols


# ---------------------------------------------------------------------------
# Chrome discovery + HTML -> PDF
# ---------------------------------------------------------------------------
CHROME_CANDIDATES = [
    "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
    "/Applications/Chromium.app/Contents/MacOS/Chromium",
    "google-chrome", "google-chrome-stable", "chromium", "chromium-browser", "chrome",
]


def find_chrome():
    for c in CHROME_CANDIDATES:
        if os.path.sep in c:
            if os.path.exists(c):
                return c
        else:
            found = shutil.which(c)
            if found:
                return found
    return None


def html_to_pdf(html_path, out_pdf, chrome_path):
    html_abs = Path(html_path).resolve()
    out_abs = Path(out_pdf).resolve()
    out_abs.parent.mkdir(parents=True, exist_ok=True)
    cmd = [
        chrome_path, "--headless", "--disable-gpu", "--no-sandbox",
        f"--print-to-pdf={out_abs}",
        "--print-to-pdf-no-header",
        "--virtual-time-budget=8000",
        "--run-all-compositor-stages-before-draw",
        f"file://{html_abs}",
    ]
    subprocess.run(cmd, check=True, capture_output=True)
    return out_abs


def report_pdf_pages(pdf_path):
    try:
        from pypdf import PdfReader
        return len(PdfReader(str(pdf_path)).pages)
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Generic fallback HTML (used for --formats pdf when no --html chart is given)
# ---------------------------------------------------------------------------
def render_fallback_html(data, title):
    parts = [
        "<!doctype html><html><head><meta charset='utf-8'>",
        f"<title>{title}</title>",
        "<style>",
        "body{font-family:-apple-system,Helvetica,Arial,sans-serif;color:#14180f;",
        "background:#fcfcfa;max-width:900px;margin:40px auto;padding:0 24px;line-height:1.5;}",
        "h1{font-size:26px;border-bottom:2px solid #14180f;padding-bottom:10px;}",
        "h2{font-size:18px;color:#1f5f5b;margin-top:32px;}",
        "h3{font-size:14px;color:#4f5348;margin-top:18px;}",
        "table{border-collapse:collapse;width:100%;font-size:12px;margin:10px 0 20px;}",
        "th,td{border:1px solid #d7dace;padding:6px 8px;text-align:left;vertical-align:top;}",
        "th{background:#eef1ea;font-family:monospace;font-size:10px;text-transform:uppercase;}",
        ".src{font-family:monospace;font-size:11px;color:#82867a;}",
        "</style></head><body>",
        f"<h1>{title}</h1>",
    ]
    if isinstance(data.get("purpose"), str):
        parts.append(f"<p>{data['purpose']}</p>")
    meta_bits = []
    if data.get("retrieved"):
        meta_bits.append(f"Retrieved: {data['retrieved']}")
    if data.get("method"):
        meta_bits.append(f"Method: {stringify_cell(data['method'], 400)}")
    if meta_bits:
        parts.append(f"<p class='src'>{' &middot; '.join(meta_bits)}</p>")

    skip = {"purpose", "retrieved", "method"}
    for key, value in data.items():
        if key in skip:
            continue
        parts.append(_html_section(key, value, level=2))
    parts.append("</body></html>")
    return "\n".join(parts)


def _html_section(key, value, level):
    tag = f"h{min(level, 6)}"
    out = [f"<{tag}>{humanize(key)}</{tag}>"]
    out.append(_html_value(value))
    return "\n".join(out)


def _html_value(value):
    if is_list_of_dicts(value):
        cols = table_columns(value)
        rows_html = ["<table><thead><tr>"] + [f"<th>{humanize(c)}</th>" for c in cols] + ["</tr></thead><tbody>"]
        for row in value:
            rows_html.append("<tr>" + "".join(f"<td>{stringify_cell(row.get(c))}</td>" for c in cols) + "</tr>")
        rows_html.append("</tbody></table>")
        return "".join(rows_html)
    if is_list_of_scalars(value):
        return "<ul>" + "".join(f"<li>{stringify_cell(v, 500)}</li>" for v in value) + "</ul>"
    if isinstance(value, list):
        return "".join(_html_value(v) if isinstance(v, (dict, list)) else f"<p>{stringify_cell(v, 500)}</p>" for v in value)
    if isinstance(value, dict):
        return "".join(_html_section(k, v, level=3) for k, v in value.items())
    return f"<p>{stringify_cell(value, 1000)}</p>"


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------
def build_docx(data, title, out_path):
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    doc.add_heading(title, level=0)
    if isinstance(data.get("purpose"), str):
        doc.add_paragraph(data["purpose"])
    meta_bits = []
    if data.get("retrieved"):
        meta_bits.append(f"Retrieved: {data['retrieved']}")
    if data.get("method"):
        meta_bits.append(f"Method: {stringify_cell(data['method'], 500)}")
    if meta_bits:
        p = doc.add_paragraph(" | ".join(meta_bits))
        for run in p.runs:
            run.italic = True
            run.font.size = Pt(9)

    skip = {"purpose", "retrieved", "method"}
    for key, value in data.items():
        if key in skip:
            continue
        _docx_section(doc, key, value, level=1)

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out_path))
    return out_path


def _docx_section(doc, key, value, level):
    doc.add_heading(humanize(key), level=min(level + 1, 9))
    _docx_value(doc, value, level)


def _docx_value(doc, value, level):
    if is_list_of_dicts(value):
        cols = table_columns(value)
        table = doc.add_table(rows=1, cols=len(cols))
        try:
            table.style = "Light Grid Accent 1"
        except KeyError:
            pass
        hdr = table.rows[0].cells
        for i, c in enumerate(cols):
            hdr[i].text = humanize(c)
        for row in value:
            cells = table.add_row().cells
            for i, c in enumerate(cols):
                cells[i].text = stringify_cell(row.get(c))
        doc.add_paragraph("")
        return
    if is_list_of_scalars(value):
        for v in value:
            doc.add_paragraph(stringify_cell(v, 1000), style="List Bullet")
        return
    if isinstance(value, list):
        for v in value:
            if isinstance(v, (dict, list)):
                _docx_value(doc, v, level + 1)
            else:
                doc.add_paragraph(stringify_cell(v, 1000))
        return
    if isinstance(value, dict):
        for k, v in value.items():
            _docx_section(doc, k, v, level + 1)
        return
    doc.add_paragraph(stringify_cell(value, 2000))


# ---------------------------------------------------------------------------
# PPTX — a summarizing deck (full detail lives in the PDF/DOCX outputs)
# ---------------------------------------------------------------------------
def flatten_for_bullets(value, depth=0, max_chars=1300):
    """Turn an arbitrary JSON value into a flat list of (indent, text) bullets,
    truncated at max_chars total so a slide's textbox stays readable."""
    lines = []
    budget = [max_chars]

    def emit(indent, text):
        if budget[0] <= 0:
            return
        text = stringify_cell(text, 220)
        lines.append((indent, text))
        budget[0] -= len(text)

    def walk(v, indent):
        if budget[0] <= 0:
            return
        if is_list_of_dicts(v):
            for row in v:
                label_key = next(iter(row.keys()), None)
                label = stringify_cell(row.get(label_key), 60) if label_key else ""
                emit(indent, label)
                for k, val in row.items():
                    if k == label_key:
                        continue
                    if isinstance(val, (dict, list)):
                        emit(indent + 1, f"{humanize(k)}: {stringify_cell(val, 150)}")
                    else:
                        emit(indent + 1, f"{humanize(k)}: {val}")
        elif is_list_of_scalars(v):
            for item in v:
                emit(indent, item)
        elif isinstance(v, list):
            for item in v:
                walk(item, indent)
        elif isinstance(v, dict):
            for k, val in v.items():
                if isinstance(val, (dict, list)):
                    emit(indent, humanize(k))
                    walk(val, indent + 1)
                else:
                    emit(indent, f"{humanize(k)}: {val}")
        else:
            emit(indent, v)

    walk(value, depth)
    if budget[0] <= 0:
        lines.append((depth, "… (truncated — see the PDF/DOCX report for full detail)"))
    return lines


def build_pptx(data, title, out_path):
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    if isinstance(data.get("purpose"), str) and len(slide.placeholders) > 1:
        slide.placeholders[1].text = stringify_cell(data["purpose"], 400)

    skip = {"purpose", "retrieved", "method"}
    for key, value in data.items():
        if key in skip:
            continue
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = humanize(key)
        body = slide.placeholders[1].text_frame
        body.word_wrap = True
        bullets = flatten_for_bullets(value)
        if not bullets:
            bullets = [(0, stringify_cell(value, 500))]
        first = True
        for indent, text in bullets:
            p = body.paragraphs[0] if first else body.add_paragraph()
            first = False
            p.text = text
            p.level = min(indent, 4)
            p.font.size = Pt(max(11, 18 - indent * 2))

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


# ---------------------------------------------------------------------------
# Orchestration
# ---------------------------------------------------------------------------
def synthesize_one(json_path, html_path, formats, out_dir, title_override=None):
    json_path = Path(json_path)
    with open(json_path) as f:
        data = json.load(f)

    title = title_override or title_from_stem(json_path.stem)
    stem = json_path.stem
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    results = {}

    if "pdf" in formats:
        try:
            if html_path and Path(html_path).exists():
                chrome = find_chrome()
                if not chrome:
                    print(f"  [pdf] SKIPPED — no Chrome/Chromium found on this machine", file=sys.stderr)
                else:
                    out_pdf = out_dir / f"{stem}.pdf"
                    html_to_pdf(html_path, out_pdf, chrome)
                    pages = report_pdf_pages(out_pdf)
                    results["pdf"] = str(out_pdf)
                    print(f"  [pdf]  {out_pdf}" + (f"  ({pages} pages)" if pages else ""))
            else:
                chrome = find_chrome()
                if not chrome:
                    print(f"  [pdf] SKIPPED — no Chrome/Chromium found on this machine", file=sys.stderr)
                else:
                    fallback_html = render_fallback_html(data, title)
                    with tempfile.NamedTemporaryFile("w", suffix=".html", delete=False) as tmp:
                        tmp.write(fallback_html)
                        tmp_path = tmp.name
                    out_pdf = out_dir / f"{stem}.pdf"
                    html_to_pdf(tmp_path, out_pdf, chrome)
                    os.unlink(tmp_path)
                    pages = report_pdf_pages(out_pdf)
                    results["pdf"] = str(out_pdf)
                    print(f"  [pdf]  {out_pdf}  (generic layout — no --html given)" + (f", {pages} pages" if pages else ""))
        except subprocess.CalledProcessError as e:
            print(f"  [pdf] FAILED: {e}", file=sys.stderr)

    if "docx" in formats:
        out_docx = out_dir / f"{stem}.docx"
        build_docx(data, title, out_docx)
        results["docx"] = str(out_docx)
        print(f"  [docx] {out_docx}")

    if "pptx" in formats:
        out_pptx = out_dir / f"{stem}.pptx"
        build_pptx(data, title, out_pptx)
        results["pptx"] = str(out_pptx)
        print(f"  [pptx] {out_pptx}")

    return results


def main():
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("json_path", nargs="?", help="Path to a bulletin's data JSON file")
    ap.add_argument("--html", help="Path to the paired chart HTML, for a pixel-faithful PDF")
    ap.add_argument("--formats", default="pdf,docx,pptx", help="Comma-separated: pdf,docx,pptx (default: all three)")
    ap.add_argument("--out-dir", default="reports", help="Output directory (default: ./reports)")
    ap.add_argument("--manifest", help="Path to repo_manifest_*.json — batch-process every entry instead of a single file")
    args = ap.parse_args()

    formats = {f.strip().lower() for f in args.formats.split(",") if f.strip()}
    unknown = formats - {"pdf", "docx", "pptx"}
    if unknown:
        ap.error(f"unknown format(s): {', '.join(unknown)} (valid: pdf, docx, pptx)")

    if args.manifest:
        manifest_path = Path(args.manifest)
        with open(manifest_path) as f:
            manifest = json.load(f)
        # Manifest entries store paths relative to the REPO ROOT (e.g. "data/x.json",
        # "charts/x.html"), matching how they're written in README.md — not relative
        # to the manifest file's own folder (data/repo_manifest_*.json).
        base = manifest_path.parent.parent
        for entry in manifest.get("entries", []):
            data_paths = entry.get("data") or []
            if not data_paths:
                print(f"\n=== {entry.get('title', '(untitled)')} — no data JSON, skipping ===")
                continue
            json_path = base / data_paths[0]
            html_path = base / entry["chart"] if entry.get("chart", "").endswith(".html") else None
            print(f"\n=== {entry.get('title', json_path.stem)} ===")
            if not json_path.exists():
                print(f"  MISSING: {json_path}", file=sys.stderr)
                continue
            synthesize_one(json_path, html_path, formats, args.out_dir, title_override=entry.get("title"))
        return

    if not args.json_path:
        ap.error("provide a json_path, or use --manifest for batch mode")

    print(f"\n=== {Path(args.json_path).stem} ===")
    synthesize_one(args.json_path, args.html, formats, args.out_dir)


if __name__ == "__main__":
    main()
